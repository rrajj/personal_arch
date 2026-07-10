import copy
from pptx import Presentation
from pptx.oxml.ns import qn

def duplicate_slide(prs, slide_index):
    """Duplicate the slide at slide_index and append it to the end. Returns new slide."""
    source = prs.slides[slide_index]
    blank_layout = source.slide_layout
    new_slide = prs.slides.add_slide(blank_layout)

    # remove any placeholder shapes add_slide() created, we'll copy the real ones
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    # copy all shapes from source slide
    for shape in source.shapes:
        new_el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(new_el)

    # copy images/media relationships (so pictures don't break)
    for rel_id, rel in source.part.rels.items():
        if "image" in rel.reltype:
            new_slide.part.rels.add_relationship(
                rel.reltype, rel.target_part, rel_id
            )

    return new_slide

def build_slide(prs, source_index, keep_shape_names):
    """Duplicate a slide, keeping only shapes whose name is in keep_shape_names."""
    new_slide = duplicate_slide(prs, source_index)
    for shape in list(new_slide.shapes):
        if shape.name not in keep_shape_names:
            shape._element.getparent().remove(shape._element)
    return new_slide


from pptx.util import Inches, Pt

def add_review_scope(prs):
    slide = basic_slide(prs, 1, keep_shapes=[_BLANK_SLIDE_SHAPES])
    set_text_by_shape_id(slide, 2, "REVIEW SCOPE")

    rows, cols = 4, 3
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4)

    graphic_frame = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = graphic_frame.table

    # header row
    table.cell(0, 0).text = "Item"
    table.cell(0, 1).text = "Owner"
    table.cell(0, 2).text = "Status"

    # data rows
    data = [
        ("Design review", "Alice", "In progress"),
        ("Code review", "Bob", "Done"),
        ("Security review", "Carol", "Pending"),
    ]
    for r, row_data in enumerate(data, start=1):
        for c, value in enumerate(row_data):
            table.cell(r, c).text = value

    return slide


from pptx.util import Pt

def style_table(table, header_size=14, body_size=11):
    # header row (row 0) — larger, bold
    for cell in table.rows[0].cells:
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(header_size)
                run.font.bold = True

    # body rows — normal size
    for row in table.rows[1:]:
        for cell in row.cells:
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(body_size)

style_table(table, header_size=14, body_size=11)


def set_cell_text(cell, text, size=12, bold=False):
    tf = cell.text_frame
    tf.clear()  # removes existing paragraphs/runs, leaves one empty paragraph
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold


