import io
from pptx import Presentation

from models import CreatePptRequest


def build_presentation(request: CreatePptRequest) -> Presentation:
    """
    Builds the PPT content from request data (architecture, presenter, pbq_answers).
    Replace with real slide-building logic pulling from temp_pts_table / temp_csi_data.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"PTS Build Review — {request.pts_id}"
    return prs


def convert_ppt_to_bytes(presentation: Presentation) -> bytes:
    """Converts an in-memory Presentation object into raw bytes for BYTEA storage."""
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()
