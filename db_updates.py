
#### 
# convert ppt to bytes
####
import io
from pptx import Presentation


def convert_ppt_to_bytes(presentation: Presentation) -> bytes:
    """
    Converts an in-memory python-pptx Presentation object into raw bytes,
    suitable for storing in a BYTEA column (ppt_blob).
    """
    buffer = io.BytesIO()
    presentation.save(buffer)
    ppt_bytes = buffer.getvalue()
    buffer.close()
    return ppt_bytes

# usage 
from pptx import Presentation

# Example: building a presentation in memory
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "PTS Build Review"

# Convert to bytes, ready for DB insert
ppt_bytes = convert_ppt_to_bytes(prs)

# Then pass into the DB function we wrote earlier
save_ppt_template(
    pts_id="PTS12345",
    permit_type="Build",
    ppt_bytes=ppt_bytes,
    file_name="PTS12345_Build_v1.pptx",
    created_by="jd123",
)



########################################################################################
# BASIC POOL AND DB ENTRY

from psycopg_pool import ConnectionPool

# Create once at app startup, reuse across requests
pool = ConnectionPool(
    conninfo="postgresql://user:password@localhost:5432/your_db",
    min_size=1,
    max_size=10,
)


def save_ppt_template(
    pts_id: str,
    permit_type: str,
    ppt_bytes: bytes,
    file_name: str,
    created_by: str,
) -> dict:
    """
    Inserts a new versioned PPT record into temp_template_store.
    Handles version increment and is_latest flip within a single transaction.
    """
    with pool.connection() as conn:
        with conn.transaction():
            with conn.cursor() as cur:
                # 1. Get next version number for this (pts_id, permit_type)
                cur.execute(
                    """
                    SELECT COALESCE(MAX(version), 0) + 1
                    FROM temp_template_store
                    WHERE pts_id = %s AND permit_type = %s
                    """,
                    (pts_id, permit_type),
                )
                next_version = cur.fetchone()[0]

                # 2. Flip old "latest" row to false, if one exists
                cur.execute(
                    """
                    UPDATE temp_template_store
                    SET is_latest = FALSE
                    WHERE pts_id = %s AND permit_type = %s AND is_latest = TRUE
                    """,
                    (pts_id, permit_type),
                )

                # 3. Insert the new version as the latest, status defaults to 'Draft'
                cur.execute(
                    """
                    INSERT INTO temp_template_store
                        (pts_id, permit_type, version, is_latest, ppt_blob,
                         file_name, file_size, created_by)
                    VALUES
                        (%s, %s, %s, TRUE, %s, %s, %s, %s)
                    RETURNING id, version, status, created_at
                    """,
                    (
                        pts_id,
                        permit_type,
                        next_version,
                        ppt_bytes,
                        file_name,
                        len(ppt_bytes),
                        created_by,
                    ),
                )
                row = cur.fetchone()

    return {
        "id": row[0],
        "version": row[1],
        "status": row[2],
        "created_at": row[3],
    }
